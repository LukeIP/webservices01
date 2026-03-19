"""
City Liveability & Urban Climate Insights API
PPTX Presentation — 5-minute oral examination slides
8 slides @ ~37s each
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Palette ──────────────────────────────────────────────────────────
NAVY  = RGBColor(0x02, 0x1B, 0x3A)
BLUE  = RGBColor(0x06, 0x5A, 0x82)
TEAL  = RGBColor(0x1C, 0x72, 0x93)
MINT  = RGBColor(0x02, 0xC3, 0x9A)
LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
PALE  = RGBColor(0xBC, 0xD9, 0xF0)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
CARD2 = RGBColor(0xE8, 0xF4, 0xFF)

# ─── Helpers ──────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=13, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    return txb

def header(slide, title):
    rect(slide, 0, 0, 10, 0.72, NAVY)
    txt(slide, title, 0.35, 0.1, 9.3, 0.55,
        size=21, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

def bg(slide, color=LIGHT):
    rect(slide, 0, 0, 10, 5.625, color)

def stat_card(slide, l, t, value, label, w=1.7, h=1.05):
    rect(slide, l, t, w, h, BLUE)
    txt(slide, value, l, t + 0.05, w, 0.55,
        size=30, color=MINT, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    txt(slide, label, l, t + 0.55, w, 0.45,
        size=9.5, color=PALE, bold=False, align=PP_ALIGN.CENTER)

def feature_card(slide, l, t, w, h, title, lines):
    rect(slide, l, t, w, h, CARD, line_color=PALE)
    rect(slide, l, t, 0.07, h, TEAL)
    txt(slide, title, l + 0.14, t + 0.06, w - 0.22, 0.28,
        size=11.5, color=BLUE, bold=True)
    y = t + 0.34
    for line in lines:
        txt(slide, f"\u2022  {line}", l + 0.14, y, w - 0.22, 0.22,
            size=9.5, color=DARK)
        y += 0.22

def layer_box(slide, l, t, w, h, label, fill, text_color=WHITE, size=11):
    rect(slide, l, t, w, h, fill)
    txt(slide, label, l, t, w, h, size=size, color=text_color,
        bold=True, align=PP_ALIGN.CENTER, wrap=False)

# ─── Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, NAVY)

# Teal left accent strip
rect(s1, 0, 0, 0.28, 5.625, TEAL)

# Mint top stripe
rect(s1, 0, 0, 10, 0.06, MINT)

# Module tag
txt(s1, "COMP  |  Web Services and Web Data  |  University of Leeds",
    0.55, 0.55, 9.0, 0.35, size=10.5, color=PALE, italic=True)

# Main title
txt(s1, "City Liveability &\nUrban Climate Insights API",
    0.55, 1.10, 8.6, 1.8, size=38, color=WHITE, bold=True, wrap=True)

# Divider line (thin rect)
rect(s1, 0.55, 2.85, 2.8, 0.04, MINT)

# Subtitle
txt(s1, "A data-driven RESTful API with MCP integration for UK city analysis",
    0.55, 3.0, 8.5, 0.5, size=14, color=PALE)

# Author + date
txt(s1, "Luke Price", 0.55, 3.8, 4.0, 0.38, size=13, color=WHITE, bold=True)
txt(s1, "March 2026", 0.55, 4.18, 4.0, 0.3, size=11, color=MUTED)

# Three stat callouts (bottom right)
stat_card(s1, 6.15, 3.55, "228", "Tests")
stat_card(s1, 7.95, 3.55, "15", "UK Cities")


# ══════════════════════════════════════════════════════════════════════
# Slide 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2)
header(s2, "Project Overview")

# Left column — description
txt(s2, "What it does", 0.4, 0.88, 5.5, 0.32, size=12, color=BLUE, bold=True)
desc_lines = [
    "Aggregates urban climate, air quality & socioeconomic data for UK cities",
    "Computes composite liveability scores across 4 weighted dimensions",
    "Detects statistical climate anomalies using z-score analysis",
    "Delivers city narratives and comparisons via structured REST endpoints",
    "Exposes an MCP server so LLM agents (e.g. Claude) can query city data conversationally",
]
y = 1.22
for line in desc_lines:
    txt(s2, f"\u2022  {line}", 0.4, y, 5.5, 0.3, size=10.5, color=DARK)
    y += 0.3

# Why it matters
txt(s2, "Why this project?", 0.4, y + 0.05, 5.5, 0.3, size=12, color=BLUE, bold=True)
txt(s2, "Sits at the intersection of environmental data and practical AI integration\n"
        "— extends beyond CRUD into an AI-accessible data service.",
    0.4, y + 0.37, 5.5, 0.55, size=10.5, color=MUTED, italic=True)

# Right column — stat cards
stat_card(s2, 6.1, 0.88, "228", "Automated Tests", w=1.75)
stat_card(s2, 7.98, 0.88, "15", "UK Cities seeded", w=1.75)
stat_card(s2, 6.1, 2.08, "9", "MCP Tools", w=1.75)
stat_card(s2, 7.98, 2.08, "4", "Analytics endpoints", w=1.75)

# Data source note
txt(s2, "Data source: Open-Meteo Archive API (free, no API key)",
    6.1, 3.25, 3.6, 0.35, size=9, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════
# Slide 3 — Architecture & Tech Stack
# ══════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3)
header(s3, "Architecture & Technology Stack")

# Left — Layered architecture diagram
txt(s3, "Application layers", 0.35, 0.85, 4.5, 0.3,
    size=11, color=BLUE, bold=True)

layers = [
    ("HTTP Clients / LLM Agents", TEAL, WHITE),
    ("FastAPI Routers", BLUE, WHITE),
    ("Service Layer (business logic)", RGBColor(0x08, 0x70, 0x9E), WHITE),
    ("SQLAlchemy ORM + Pydantic schemas", RGBColor(0x0A, 0x85, 0xBB), WHITE),
    ("SQLite (dev)  /  PostgreSQL (prod)", NAVY, WHITE),
]
y_l = 1.2
for label, fill, tc in layers:
    layer_box(s3, 0.35, y_l, 4.5, 0.53, label, fill, tc, size=10)
    if y_l < 3.3:  # arrow
        txt(s3, "\u25bc", 2.35, y_l + 0.53, 0.8, 0.18, size=10, color=TEAL,
            align=PP_ALIGN.CENTER)
    y_l += 0.71

# Right — Tech stack table
txt(s3, "Technology choices", 5.2, 0.85, 4.4, 0.3,
    size=11, color=BLUE, bold=True)

stack = [
    ("Framework",  "FastAPI — async, auto OpenAPI docs, Pydantic v2"),
    ("ORM",        "SQLAlchemy 2.0 — backend-agnostic, mature query API"),
    ("Database",   "SQLite (dev) / PostgreSQL via Railway (prod)"),
    ("Migrations", "Alembic — versioned, rollback-safe schema changes"),
    ("Auth",       "JWT (HS256) + bcrypt hashing + RBAC (user / admin)"),
    ("MCP",        "fastmcp SDK — 9 tools, mounted at /mcp/sse"),
    ("Ext. data",  "Open-Meteo — free historical weather, no API key"),
]
y_r = 1.2
for label, desc in stack:
    rect(s3, 5.2, y_r, 4.4, 0.44, CARD, line_color=PALE)
    txt(s3, label, 5.28, y_r + 0.04, 1.1, 0.38, size=9.5, color=BLUE, bold=True)
    txt(s3, desc, 6.42, y_r + 0.04, 3.1, 0.38, size=9, color=DARK)
    y_r += 0.48


# ══════════════════════════════════════════════════════════════════════
# Slide 4 — API Features & Endpoints
# ══════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4)
header(s4, "API Features & Key Endpoints")

# 4 feature cards in 2x2 grid
feature_card(s4, 0.3, 0.85, 4.6, 1.85, "CRUD Operations", [
    "Cities — POST / GET / PUT / DELETE /api/v1/cities",
    "Users — register, update profile, admin management",
    "Climate readings — ingest, update, bulk delete",
    "Auto-seeded with 365 days of real Open-Meteo data",
])
feature_card(s4, 5.1, 0.85, 4.6, 1.85, "Analytics", [
    "GET /cities/{id}/liveability — composite score (4 dims)",
    "GET /cities/compare — side-by-side city ranking",
    "GET /cities/{id}/anomalies — z-score outlier detection",
    "GET /cities/{id}/trends — historical climate trends",
])
feature_card(s4, 0.3, 2.85, 4.6, 1.65, "Security", [
    "JWT stateless auth — POST /auth/login",
    "RBAC — user vs. admin roles, injected at router level",
    "Rate limiting — 100 req/min via slowapi middleware",
    "bcrypt password hashing, tamper-proof token validation",
])
feature_card(s4, 5.1, 2.85, 4.6, 1.65, "MCP Server (Advanced)", [
    "9 AI tools: search, liveability, anomalies, trends, record",
    "Mounted at /mcp/sse — no separate deployment needed",
    "Claude Desktop connects to Railway URL directly",
    "LLMs query UK city data conversationally via tool calls",
])


# ══════════════════════════════════════════════════════════════════════
# Slide 5 — Version Control & Testing
# ══════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5)
header(s5, "Version Control & Testing")

# Left — Version control
rect(s5, 0.3, 0.85, 4.5, 0.38, BLUE)
txt(s5, "Version Control Practices", 0.4, 0.88, 4.3, 0.32,
    size=12, color=WHITE, bold=True)

vc_items = [
    ("Public GitHub repo", "Visible commit history from initial scaffold to final submission"),
    ("Structured commits", "Descriptive messages — feature, fix, test, config scopes"),
    ("Alembic migrations", "All schema changes versioned and reversible in the repo"),
    ("README.md", "Setup instructions, environment variables, local run guide"),
    ("Environment config", "SKIP_SEED flag prevents external HTTP calls during testing"),
]
y = 1.3
for heading, detail in vc_items:
    txt(s5, heading, 0.4, y, 1.65, 0.22, size=10, color=TEAL, bold=True)
    txt(s5, detail,  2.1, y, 2.6, 0.22, size=9.5, color=DARK)
    y += 0.26

# Right — Testing
rect(s5, 5.2, 0.85, 4.5, 0.38, TEAL)
txt(s5, "228 Tests across 11 modules", 5.3, 0.88, 4.3, 0.32,
    size=12, color=WHITE, bold=True)

test_rows = [
    ("Integration", "Full req/resp cycles, in-memory SQLite, seeded fixtures"),
    ("Unit",        "Service layer with real (temporary) SQLite sessions"),
    ("Security",    "401 on unauth, 403 on non-admin, tampered JWT rejection"),
    ("Utility",     "Scoring algorithm & SQL validator with edge-case inputs"),
    ("Middleware",  "Rate limiter triggers at threshold; latency logging fields"),
]
y = 1.3
for label, desc in test_rows:
    rect(s5, 5.2, y, 1.05, 0.24, BLUE)
    txt(s5, label, 5.22, y + 0.02, 1.01, 0.21, size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s5, desc, 6.3, y, 3.3, 0.24, size=9.5, color=DARK)
    y += 0.28

# Key decision note
txt(s5, "Key decision: real SQLite sessions over mocked DB calls — mocks can mask "
        "schema-level bugs that only appear when the ORM generates actual SQL.",
    5.2, 2.74, 4.5, 0.55, size=9, color=MUTED, italic=True)

# Divider
rect(s5, 4.85, 0.85, 0.04, 4.4, PALE)


# ══════════════════════════════════════════════════════════════════════
# Slide 6 — API Documentation
# ══════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6)
header(s6, "API Documentation")

# Three panels
panels = [
    ("Swagger UI",  "/docs",
     ["Interactive try-it-out for every endpoint",
      "Auto-generated from Python type annotations",
      "Auth flow: Authorize button with JWT bearer"],
     BLUE),
    ("ReDoc",       "/redoc",
     ["Clean, reference-style documentation",
      "Schema definitions with example payloads",
      "Accessible without authentication"],
     TEAL),
    ("PDF Export",  "docs/api_documentation.pdf",
     ["Submitted to Minerva and linked in README.md",
      "Generated from Swagger JSON",
      "Full endpoint list with request/response schemas"],
     RGBColor(0x02, 0x84, 0x6A)),
]
x = 0.3
for title, path, bullets, color in panels:
    rect(s6, x, 0.85, 3.1, 0.42, color)
    txt(s6, title, x + 0.1, 0.88, 2.9, 0.36, size=13, color=WHITE, bold=True)
    rect(s6, x, 1.27, 3.1, 0.3, DARK)
    txt(s6, path, x + 0.1, 1.3, 2.9, 0.25, size=9.5, color=MINT, italic=True)
    y = 1.63
    for b in bullets:
        txt(s6, f"\u2022  {b}", x + 0.1, y, 2.9, 0.25, size=9.5, color=DARK)
        y += 0.3
    x += 3.25

# What is documented
txt(s6, "All endpoints document:", 0.3, 3.65, 4.5, 0.3,
    size=11, color=BLUE, bold=True)
doc_fields = [
    "HTTP method, path, and path/query parameters",
    "Request body schema (Pydantic models)",
    "Response model and status codes (200, 201, 401, 403, 404, 422, 429)",
    "Authentication requirements and example payloads",
]
y = 3.98
for f in doc_fields:
    txt(s6, f"\u2022  {f}", 0.3, y, 4.5, 0.24, size=10, color=DARK)
    y += 0.24

# Compliance note
rect(s6, 5.1, 3.65, 4.6, 0.98, CARD2, line_color=PALE)
txt(s6, "Submission compliance", 5.22, 3.7, 4.36, 0.28,
    size=10.5, color=BLUE, bold=True)
txt(s6, "PDF referenced in README.md\n"
        "Hosted in GitHub repo under /docs/\n"
        "Also linked in Technical Report appendix",
    5.22, 3.98, 4.36, 0.6, size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════
# Slide 7 — Technical Report Highlights
# ══════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg(s7)
header(s7, "Technical Report Highlights")

highlights = [
    ("Stack Justification",
     "FastAPI chosen over Django for native async, auto-OpenAPI, and Pydantic v2 throughput. "
     "NoSQL ruled out — relational JOINs across cities, climate, and socioeconomic data are more efficient."),
    ("Scoring Model",
     "Composite liveability: climate comfort 25%, affordability 30%, safety 25%, environment 20%. "
     "Each metric normalised 0–100 using domain-specific functions (e.g. AQI inverted)."),
    ("Anomaly Detection",
     "Z-score method: compute mean/std across all stored readings for a city+metric pair; flag "
     "dates where |z| > 2.0. Statistically sound; no external ML library required."),
    ("Key Design Challenge",
     "Router ordering — GET /cities/compare must register before GET /cities/{city_id} or "
     "FastAPI matches 'compare' as a city_id. Documented in main.py."),
    ("GenAI Usage",
     "AI used at ideation, scaffold generation, test-case review, and report drafting. Freed time "
     "for high-level decisions (scoring model, MCP design). Limitation: no proactive refactoring."),
]
y = 0.88
for title, body in highlights:
    rect(s7, 0.3, y, 9.4, 0.72, CARD, line_color=PALE)
    rect(s7, 0.3, y, 0.07, 0.72, MINT)
    txt(s7, title, 0.48, y + 0.04, 2.1, 0.3, size=10.5, color=BLUE, bold=True)
    txt(s7, body, 2.65, y + 0.06, 6.95, 0.6, size=9.5, color=DARK)
    y += 0.8


# ══════════════════════════════════════════════════════════════════════
# Slide 8 — Deliverables (dark closing)
# ══════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
bg(s8, NAVY)
rect(s8, 0, 0, 10, 0.06, MINT)

txt(s8, "Deliverables", 0.5, 0.22, 9.0, 0.55,
    size=26, color=WHITE, bold=True)

deliverables = [
    ("Code Repository",
     "Public GitHub — versioned source, commit history, README.md with setup guide"),
    ("API Documentation",
     "Swagger UI (/docs), ReDoc (/redoc), PDF in /docs/ — all endpoints, auth, error codes"),
    ("Technical Report",
     "5-page PDF — stack justification, architecture, testing, GenAI declaration (Appendix)"),
    ("Presentation Slides",
     "This PPTX — version control, API docs, report highlights, all deliverables"),
]

y = 0.95
for i, (title, desc) in enumerate(deliverables):
    # Number badge
    rect(s8, 0.4, y, 0.5, 0.5, TEAL)
    txt(s8, str(i + 1), 0.4, y, 0.5, 0.5, size=16, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    # Card
    rect(s8, 1.05, y, 8.5, 0.5, RGBColor(0x05, 0x2D, 0x5E))
    txt(s8, title, 1.18, y + 0.04, 2.1, 0.28, size=11, color=MINT, bold=True)
    txt(s8, desc,  3.38, y + 0.06, 6.1, 0.38, size=10, color=PALE)
    y += 0.62

# Closing note
txt(s8, "Live demo available — Railway deployment at API base URL",
    0.4, 3.58, 9.2, 0.38, size=12, color=TEAL, bold=True,
    align=PP_ALIGN.CENTER)
txt(s8, "Questions welcome on any aspect of the implementation, design decisions, or GenAI usage",
    0.4, 3.98, 9.2, 0.36, size=11, color=PALE, align=PP_ALIGN.CENTER)

# Footer
rect(s8, 0, 5.3, 10, 0.325, RGBColor(0x01, 0x10, 0x22))
txt(s8, "Luke Price  |  COMP Web Services  |  University of Leeds  |  March 2026",
    0, 5.3, 10, 0.325, size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ─── Save ──────────────────────────────────────────────────────────────
out = r"C:\Users\lukep\Desktop\webservices01\docs\presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
